# =======
# Imports
# =======
from __future__ import annotations  # Enable postponed type hints.
import io  # Handle in-memory byte streams.
import os  # Handle path utilities.
import zipfile  # Inspect Office archive files.
import mammoth  # Convert DOCX content to HTML.
import nh3  # Sanitize extracted HTML.
from html import escape  # Escape unsafe HTML text.
from typing import Any  # Support flexible type hints.
from bs4 import BeautifulSoup, NavigableString, Tag  # Parse extracted HTML content.
from mammoth.images import img_element  # Convert DOCX images during parsing.
from PIL import Image, UnidentifiedImageError  # Normalise extracted image formats.
from pptx import Presentation  # Read PowerPoint presentations.
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER  # Inspect PowerPoint shape types.

# ======
# Limits
# ======
SUPPORTED_EXTENSIONS = {".docx", ".pptx"}  # Define allowed Office extensions.
MAX_PARSED_UPLOAD_BYTES = 50 * 1024 * 1024  # Set the upload size limit.
MAX_OFFICE_ARCHIVE_MEMBERS = 2000  # Limit archive member count.
MAX_OFFICE_UNCOMPRESSED_BYTES = 100 * 1024 * 1024  # Limit expanded archive size.

# ==============
# Sanitised HTML
# ==============
SAFE_TAGS = {  # Define allowed HTML tags.
    "a",  # Allow anchor tags.
    "blockquote",  # Allow blockquote tags.
    "br",  # Allow line breaks.
    "em",  # Allow emphasis tags.
    "figcaption",  # Allow figure captions.
    "figure",  # Allow figure wrappers.
    "h1",  # Allow level one headings.
    "h2",  # Allow level two headings.
    "h3",  # Allow level three headings.
    "h4",  # Allow level four headings.
    "h5",  # Allow level five headings.
    "h6",  # Allow level six headings.
    "hr",  # Allow horizontal rules.
    "img",  # Allow image tags.
    "li",  # Allow list items.
    "ol",  # Allow ordered lists.
    "p",  # Allow paragraph tags.
    "strong",  # Allow strong emphasis.
    "table",  # Allow table tags.
    "tbody",  # Allow table bodies.
    "td",  # Allow table cells.
    "th",  # Allow table headers.
    "thead",  # Allow table headers.
    "tr",  # Allow table rows.
    "ul",  # Allow unordered lists.
}

SAFE_ATTRIBUTES = {  # Define allowed HTML attributes.
    "a": {"href", "target"},  # Allow anchor attributes.
    "img": {"src", "alt"},  # Allow image attributes.
    "th": {"colspan", "rowspan"},  # Allow table-header attributes.
    "td": {"colspan", "rowspan"},  # Allow table-cell attributes.
}

# ==============
# Public Parsing
# ==============
def validate_supported_upload(uploaded_file) -> str:  # Define the upload validator.
    """Validate the uploaded file."""
    name = getattr(uploaded_file, "name", "") or ""  # Read the uploaded filename.
    _, ext = os.path.splitext(name)  # Split the filename extension.
    ext = ext.lower()  # Extract the file extension.

    size = getattr(uploaded_file, "size", 0) or 0  # Read the uploaded file size.

    if ext not in SUPPORTED_EXTENSIONS:  # Check the current condition.
        raise ValueError(  # Raise a validation error.
            "Only .docx and .pptx files are allowed for weekly notes and lecturer assignment materials."  # Explain the allowed file types.
        )

    if size <= 0:  # Check the current condition.
        raise ValueError("The uploaded file is empty.")  # Raise a validation error.

    if size > MAX_PARSED_UPLOAD_BYTES:  # Check the current condition.
        raise ValueError("The uploaded file exceeds the 15 MB limit.")  # Raise a validation error.

    return ext  # Return the computed value.

def _validate_office_container(file_bytes: bytes, extension: str) -> None:  # Define the archive validator.
    """Validate the Office archive container."""
    try:  # Start guarded parsing.
        with zipfile.ZipFile(io.BytesIO(file_bytes)) as archive:  # Open the resource safely.
            infos = archive.infolist()  # Read archive member metadata.
            names = set(archive.namelist())  # Collect archive member names.

            if len(infos) > MAX_OFFICE_ARCHIVE_MEMBERS:  # Check the current condition.
                raise ValueError("The uploaded file is too complex to process safely.")  # Raise a validation error.

            total_uncompressed = sum(info.file_size for info in infos)  # Sum expanded archive size.
            if total_uncompressed > MAX_OFFICE_UNCOMPRESSED_BYTES:  # Check the current condition.
                raise ValueError("The uploaded file expands beyond the safe processing limit.")  # Raise a validation error.

    except zipfile.BadZipFile as exc:  # Handle parsing failures.
        raise ValueError("The uploaded file is not a valid Office document.") from exc  # Raise a validation error.

    required_entries = {  # Define required Office entries.
        ".docx": {"[Content_Types].xml", "word/document.xml"},  # Require core DOCX entries.
        ".pptx": {"[Content_Types].xml", "ppt/presentation.xml"},  # Require core PPTX entries.
    }[extension]  # Select entries for this extension.

    if not required_entries.issubset(names):  # Check the current condition.
        raise ValueError("The uploaded file is not a valid Office document.")  # Raise a validation error.

def parse_uploaded_office_file(uploaded_file) -> dict[str, Any]:  # Define the upload parser.
    """Parse an uploaded Office file."""
    extension = validate_supported_upload(uploaded_file)  # Store the detected extension.
    file_bytes = uploaded_file.read()  # Read the uploaded file bytes.

    if hasattr(uploaded_file, "seek"):  # Check the current condition.
        uploaded_file.seek(0)  # Reset the file pointer.

    if not file_bytes:  # Check the current condition.
        raise ValueError("The uploaded file is empty.")  # Raise a validation error.


    _validate_office_container(file_bytes, extension)  # Validate the Office archive.

    if extension == ".docx":  # Check the current condition.
        parsed = parse_docx_file(file_bytes)  # Store the parsed document payload.

    else:  # Handle the fallback case.
        parsed = parse_pptx_file(file_bytes)  # Store the parsed document payload.

    parsed["extension"] = extension.lstrip(".")  # Store the extension without the dot.
    return parsed  # Return the computed value.

# ============
# DOCX Helpers
# ============
def parse_docx_file(file_bytes: bytes) -> dict[str, Any]:  # Define the DOCX parser.
    """Parse a DOCX file."""
    captured_images: list[dict[str, Any]] = []  # Initialise captured images.
    image_counter = 0  # Track extracted image order.

    def convert_image(image):  # Define convert image.
        """Capture a DOCX image."""
        nonlocal image_counter  # Reuse the outer counter.

        image_counter += 1  # Increment the image counter.
        token = f"img-{image_counter}"  # Build a stable image token.

        extension = _extension_from_content_type(getattr(image, "content_type", "")) or "png"  # Store the detected extension.
        filename = f"{token}.{extension}"  # Build the stored image name.

        with image.open() as image_bytes:  # Open the resource safely.
            blob = image_bytes.read()  # Read the raw image bytes.

        blob, extension, filename = _normalise_image_blob(blob, extension, filename)  # Normalise the extracted image.

        captured_images.append(  # Store the extracted image.
            {  # Build the payload dictionary.
                "token": token,  # Store the image token.
                "filename": filename,  # Store the image filename.
                "content": blob,  # Store raw image bytes.
                "page_number": None,  # Leave the page unset.
                "alt_text": getattr(image, "alt_text", "") or "",  # Store the image alt text.
                "display_order": image_counter,  # Store the image order.
            }
        )

        return {  # Return the parsed payload.
            "src": f"parsed-image:{token}",  # Point to the parsed image.
            "alt": getattr(image, "alt_text", "") or "",  # Store the image alt text.
        }

    result = mammoth.convert_to_html(  # Store Mammoth conversion output.
        io.BytesIO(file_bytes),  # Wrap bytes for Mammoth.
        convert_image=img_element(convert_image),  # Set convert image.
        include_embedded_style_map=False,  # Set include embedded style map.
    )

    blocks = _docx_html_to_blocks(result.value)  # Store parsed content blocks.

    return {  # Return the parsed payload.
        "blocks": blocks,  # Return parsed content blocks.
        "page_count": len(blocks),  # Return the page count.
        "images": captured_images,  # Return extracted images.
        "warnings": [str(message) for message in getattr(result, "messages", [])],  # Return Mammoth warnings.
    }

# ============
# PPTX Helpers
# ============
def parse_pptx_file(file_bytes: bytes) -> dict[str, Any]:  # Define the PPTX parser.
    """Parse a PPTX file."""
    presentation = Presentation(io.BytesIO(file_bytes))  # Open the PowerPoint file.

    blocks: list[dict[str, Any]] = []  # Initialise parsed blocks.
    captured_images: list[dict[str, Any]] = []  # Initialise captured images.
    image_counter = 0  # Track extracted image order.

    for slide_index, slide in enumerate(presentation.slides, start=1):  # Iterate through the collection.
        page_elements: list[dict[str, Any]] = []  # Initialise slide elements.

        slide_title = ""  # Store the slide title.
        title_shape = getattr(slide.shapes, "title", None)  # Read the title placeholder.
        if title_shape and getattr(title_shape, "text", "").strip():  # Check the current condition.
            slide_title = title_shape.text.strip()  # Store the slide title.

        for shape in _sorted_shapes(slide.shapes):  # Iterate through the collection.
            extracted_elements, image_counter, new_images = _extract_pptx_shape_content(  # Extract content from the shape.
                shape=shape,  # Set shape.
                slide_index=slide_index,  # Set slide index.
                image_counter=image_counter,  # Track extracted image order.
            )
            if extracted_elements:  # Check the current condition.
                page_elements.extend(extracted_elements)  # Append extracted elements.
            if new_images:  # Check the current condition.
                captured_images.extend(new_images)  # Append extracted images.

        if not page_elements:  # Check the current condition.
            continue  # Skip to the next item.

        blocks.append(  # Append the page block.
            {  # Build the payload dictionary.
                "type": "page",  # Mark this block as a page.
                "page_number": slide_index,  # Store the slide number.
                "label": slide_title or f"Slide {slide_index}",  # Use the slide label.
                "elements": page_elements,  # Store page elements.
            }
        )

    return {  # Return the parsed payload.
        "blocks": blocks,  # Return parsed content blocks.
        "page_count": len(blocks),  # Return the page count.
        "images": captured_images,  # Return extracted images.
        "warnings": [],  # Return an empty warning list.
    }

def build_rendered_html_from_blocks(  # Define the HTML renderer.
    blocks: list[dict[str, Any]],  # Initialise parsed blocks.
    image_lookup: dict[str, dict[str, str]],  # Annotate stored images.
) -> str:  # Finish the function signature.
    """Build rendered HTML from blocks."""
    page_html: list[str] = []  # Annotate page html.
    total_pages = len(blocks)  # Set total pages.

    for index, page in enumerate(blocks, start=1):  # Iterate through the collection.
        label = escape(str(page.get("label") or f"Page {index}"))  # Build the page label.
        inner_html: list[str] = []  # Annotate inner html.

        for element in page.get("elements", []):  # Iterate through the collection.
            if element.get("type") != "raw_html":  # Check the current condition.
                continue  # Skip to the next item.

            snippet = _hydrate_image_placeholders(  # Store the current HTML snippet.
                element.get("html", ""),  # Read the raw HTML snippet.
                image_lookup=image_lookup,  # Set image lookup.
            )
            snippet = _sanitize_user_html(snippet)  # Store the current HTML snippet.

            if "<table" in snippet:  # Check the current condition.
                snippet = f'<div class="parsed-table-wrap">{snippet}</div>'  # Store the current HTML snippet.

            inner_html.append(snippet)  # Store the rendered snippet.

        page_html.append(  # Append the page HTML.
            f'<section class="parsed-page" aria-label="{label}">'  # Build this path segment.
            f'<div class="parsed-page-header">{label}</div>'  # Build this path segment.
            f'{"".join(inner_html)}'  # Render the inner page HTML.
            f"</section>"  # Close the page section.
        )

        if index < total_pages:  # Check the current condition.
            page_html.append('<hr class="parsed-page-break" aria-hidden="true">')  # Insert a page break.

    return "".join(page_html)  # Return the computed value.

def _docx_html_to_blocks(html: str) -> list[dict[str, Any]]:  # Define the DOCX block builder.
    """Convert DOCX HTML into blocks."""
    soup = BeautifulSoup(html or "", "html.parser")  # Parse the HTML fragment.
    body_nodes: list[Tag] = []  # Initialise body nodes.

    for child in soup.contents:  # Iterate through the collection.
        if isinstance(child, NavigableString):  # Check the current condition.
            if child.strip():  # Check the current condition.
                frag = BeautifulSoup(f"<p>{escape(str(child).strip())}</p>", "html.parser")  # Wrap loose text in a paragraph.
                if frag.p:  # Check the current condition.
                    body_nodes.append(frag.p)  # Store the wrapped node.
            continue  # Skip to the next item.

        if isinstance(child, Tag):  # Check the current condition.
            body_nodes.append(child)  # Store the HTML node.

    pages: list[dict[str, Any]] = []  # Initialise page blocks.
    current_page = {  # Track the current page block.
        "type": "page",  # Mark this block as a page.
        "page_number": 1,  # Store the first page number.
        "label": "Page 1",  # Store the first page label.
        "elements": [],  # Start with no page elements.
    }

    for node in body_nodes:  # Iterate through the collection.
        if node.name in {"h1", "h2"} and current_page["elements"]:  # Check the current condition.
            pages.append(current_page)  # Store the completed page.
            next_number = len(pages) + 1  # Calculate the next page number.
            current_page = {  # Track the current page block.
                "type": "page",  # Mark this block as a page.
                "page_number": next_number,  # Store the next page number.
                "label": f"Page {next_number}",  # Store the next page label.
                "elements": [],  # Start with no page elements.
            }

        snippet = str(node)  # Store the current HTML snippet.
        if snippet.strip():  # Check the current condition.
            current_page["elements"].append(  # Append the page element.
                {  # Build the payload dictionary.
                    "type": "raw_html",  # Mark this element as raw HTML.
                    "html": snippet,  # Store the raw HTML.
                }
            )

    if current_page["elements"]:  # Check the current condition.
        pages.append(current_page)  # Store the completed page.

    if not pages:  # Check the current condition.
        pages.append(  # Append the fallback page.
            {  # Build the payload dictionary.
                "type": "page",  # Mark this block as a page.
                "page_number": 1,  # Store the first page number.
                "label": "Page 1",  # Store the first page label.
                "elements": [  # Store this dictionary value.
                    {  # Build the payload dictionary.
                        "type": "raw_html",  # Mark this element as raw HTML.
                        "html": "<p>No readable content could be extracted from this document.</p>",  # Store this dictionary value.
                    }
                ],
            }
        )

    return pages  # Return the computed value.

def _sorted_shapes(shapes) -> list[Any]:  # Define the shape sorter.
    """Sort PowerPoint shapes."""
    ordered: list[Any] = []  # Initialise sorted shapes.

    for shape in shapes:  # Iterate through the collection.
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:  # Check the current condition.
            ordered.extend(_sorted_shapes(shape.shapes))  # Flatten grouped shapes.
        else:  # Handle the fallback case.
            ordered.append(shape)  # Store the current shape.

    return sorted(ordered, key=lambda shp: (getattr(shp, "top", 0), getattr(shp, "left", 0)))  # Return the computed value.

def _extract_pptx_shape_content(shape, slide_index: int, image_counter: int):  # Define the shape extractor.
    """Extract slide shape content."""
    elements: list[dict[str, Any]] = []  # Initialise extracted elements.
    images: list[dict[str, Any]] = []  # Initialise extracted images.

    shape_type = getattr(shape, "shape_type", None)  # Read the current shape type.

    if shape_type == MSO_SHAPE_TYPE.PICTURE:  # Check the current condition.
        image_counter += 1  # Increment the image counter.
        token = f"img-{image_counter}"  # Build a stable image token.

        img = shape.image  # Read the picture payload.
        extension = (getattr(img, "ext", "") or "png").lower()  # Store the detected extension.
        filename = getattr(img, "filename", "") or f"{token}.{extension}"  # Build the stored image name.
        blob = img.blob  # Read the raw image bytes.

        blob, extension, filename = _normalise_image_blob(blob, extension, filename)  # Normalise the extracted image.
        alt_text = _best_picture_alt_text(shape, slide_index)  # Build the image alt text.

        images.append(  # Store the picture metadata.
            {  # Build the payload dictionary.
                "token": token,  # Store the image token.
                "filename": filename,  # Store the image filename.
                "content": blob,  # Store raw image bytes.
                "page_number": slide_index,  # Store the slide number.
                "alt_text": alt_text,  # Store the image alt text.
                "display_order": image_counter,  # Store the image order.
            }
        )

        elements.append(  # Store the rendered element.
            {  # Build the payload dictionary.
                "type": "raw_html",  # Mark this element as raw HTML.
                "html": (  # Build the figure HTML.
                    f'<figure><img src="parsed-image:{token}" alt="{escape(alt_text)}"></figure>'  # Build this path segment.
                ),
            }
        )
        return elements, image_counter, images  # Return the computed value.

    if getattr(shape, "has_table", False):  # Check the current condition.
        table_html = _table_to_html(shape.table)  # Render the table as HTML.
        if table_html:  # Check the current condition.
            elements.append({"type": "raw_html", "html": table_html})  # Store the table element.
        return elements, image_counter, images  # Return the computed value.

    if getattr(shape, "has_text_frame", False):  # Check the current condition.
        text_html = _text_frame_to_html(shape)  # Render the text frame as HTML.
        if text_html:  # Check the current condition.
            elements.append({"type": "raw_html", "html": text_html})  # Store the text element.
        return elements, image_counter, images  # Return the computed value.

    return elements, image_counter, images  # Return the computed value.

def _text_frame_to_html(shape) -> str:  # Define the text-frame renderer.
    """Convert a text frame to HTML."""
    text_frame = shape.text_frame  # Read the text frame.
    paragraphs = [p for p in text_frame.paragraphs if (p.text or "").strip()]  # Drop the handled paragraph.

    if not paragraphs:  # Check the current condition.
        return ""  # Return the computed value.

    placeholder_type = None  # Reset the placeholder type.
    if getattr(shape, "is_placeholder", False):  # Check the current condition.
        try:  # Start guarded parsing.
            placeholder_type = shape.placeholder_format.type  # Track the placeholder type.
        except Exception:  # Handle parsing failures.
            placeholder_type = None  # Reset the placeholder type.

    html_parts: list[str] = []  # Initialise HTML parts.

    if placeholder_type in {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}:  # Check the current condition.
        html_parts.append(f"<h2>{_paragraph_inline_html(paragraphs[0])}</h2>")  # Render the title heading.
        paragraphs = paragraphs[1:]  # Drop the handled paragraph.
    elif placeholder_type == PP_PLACEHOLDER.SUBTITLE:  # Check the next condition.
        html_parts.append(f"<h3>{_paragraph_inline_html(paragraphs[0])}</h3>")  # Render the subtitle heading.
        paragraphs = paragraphs[1:]  # Drop the handled paragraph.

    if not paragraphs:  # Check the current condition.
        return "".join(html_parts)  # Return the computed value.

    should_render_as_list = (  # Decide whether to render a list.
        placeholder_type in {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT}  # Check body-style placeholders.
        and len(paragraphs) > 1  # Require multiple paragraphs.
    )

    if should_render_as_list:  # Check the current condition.
        html_parts.append("<ul>")  # Open the list wrapper.
        for paragraph in paragraphs:  # Iterate through the collection.
            html_parts.append(f"<li>{_paragraph_inline_html(paragraph)}</li>")  # Render the list item.
        html_parts.append("</ul>")  # Close the list wrapper.
    else:  # Handle the fallback case.
        for paragraph in paragraphs:  # Iterate through the collection.
            html_parts.append(f"<p>{_paragraph_inline_html(paragraph)}</p>")  # Render the paragraph.

    return "".join(html_parts)  # Return the computed value.

def _paragraph_inline_html(paragraph) -> str:  # Define the inline HTML builder.
    """Build inline HTML for a paragraph."""
    fragments: list[str] = []  # Initialise inline fragments.
    runs = list(getattr(paragraph, "runs", []))  # Collect text runs.

    if not runs:  # Check the current condition.
        return escape((paragraph.text or "").strip())  # Return the computed value.

    for run in runs:  # Iterate through the collection.
        text = escape(run.text or "")  # Escape the run text.
        if not text:  # Check the current condition.
            continue  # Skip to the next item.

        href = getattr(getattr(run, "hyperlink", None), "address", None)  # Read the hyperlink target.
        if href:  # Check the current condition.
            text = (  # Escape the run text.
                f'<a href="{escape(href, quote=True)}" '  # Build this path segment.
                f'target="_blank">{text}</a>'  # Build this path segment.
            )

        if getattr(getattr(run, "font", None), "bold", False):  # Check the current condition.
            text = f"<strong>{text}</strong>"  # Escape the run text.

        if getattr(getattr(run, "font", None), "italic", False):  # Check the current condition.
            text = f"<em>{text}</em>"  # Escape the run text.

        fragments.append(text)  # Store the rendered fragment.

    return "".join(fragments).strip()  # Return the computed value.

def _table_to_html(table) -> str:  # Define the table renderer.
    """Convert a table to HTML."""
    rows_html: list[str] = []  # Initialise table rows.

    for row_index, row in enumerate(table.rows):  # Iterate through the collection.
        cell_tag = "th" if row_index == 0 else "td"  # Choose the cell tag.
        cell_html: list[str] = []  # Initialise cell HTML.

        for cell in row.cells:  # Iterate through the collection.
            text = escape(cell.text.strip()) if cell.text else ""  # Escape the run text.
            cell_html.append(f"<{cell_tag}>{text}</{cell_tag}>")  # Render the table cell.

        rows_html.append(f"<tr>{''.join(cell_html)}</tr>")  # Render the table row.

    if not rows_html:  # Check the current condition.
        return ""  # Return the computed value.

    return f"<table><tbody>{''.join(rows_html)}</tbody></table>"  # Return the display string.

def _hydrate_image_placeholders(  # Define the image hydrator.
    snippet: str,  # Annotate snippet.
    image_lookup: dict[str, dict[str, str]],  # Annotate stored images.
) -> str:  # Finish the function signature.
    """Replace parsed image placeholders."""
    soup = BeautifulSoup(snippet or "", "html.parser")  # Parse the HTML fragment.

    for img in soup.find_all("img"):  # Iterate through the collection.
        src = (img.get("src") or "").strip()  # Read the image source.

        if not src.startswith("parsed-image:"):  # Check the current condition.
            continue  # Skip to the next item.

        token = src.split("parsed-image:", 1)[1]  # Build a stable image token.
        data = image_lookup.get(token)  # Read the image lookup data.

        if not data:  # Check the current condition.
            img.decompose()  # Remove the missing image.
            continue  # Skip to the next item.

        img["src"] = data["src"]  # Replace the image source.
        img["alt"] = data.get("alt_text", "") or img.get("alt", "") or ""  # Replace the image alt text.

    return str(soup)  # Return the computed value.

def _sanitize_user_html(html: str) -> str:  # Define the HTML sanitizer.
    """Sanitize extracted user HTML."""
    return nh3.clean(  # Return the computed value.
        html or "",  # Clean the provided HTML.
        tags=SAFE_TAGS,  # Set tags.
        attributes=SAFE_ATTRIBUTES,  # Set attributes.
        url_schemes={"http", "https", "mailto"},  # Set url schemes.
        link_rel="noopener noreferrer",  # Set link rel.
    )

def _best_picture_alt_text(shape, slide_index: int) -> str:  # Define the alt-text helper.
    """Get picture alt text."""
    try:  # Start guarded parsing.
        descr = shape._element.xpath(  # Read the image description.
            ".//p:cNvPr/@descr",  # Allow the .//p:cNvPr/@descr value.
            namespaces={"p": "http://schemas.openxmlformats.org/presentationml/2006/main"},  # Set namespaces.
        )
        if descr and descr[0].strip():  # Check the current condition.
            return descr[0].strip()  # Return the computed value.
    except Exception:  # Handle parsing failures.
        pass  # Ignore the failure safely.

    return f"Image from slide {slide_index}"  # Return the display string.

# =============
# Image Helpers
# =============
def _extension_from_content_type(content_type: str) -> str:  # Define the extension mapper.
    """Map content type to an extension."""
    mapping = {  # Map MIME types to extensions.
        "image/jpeg": "jpg",  # Define a MIME mapping.
        "image/png": "png",  # Define a MIME mapping.
        "image/gif": "gif",  # Define a MIME mapping.
        "image/svg+xml": "svg",  # Define a MIME mapping.
        "image/webp": "webp",  # Define a MIME mapping.
        "image/tiff": "tiff",  # Define a MIME mapping.
        "image/bmp": "bmp",  # Define a MIME mapping.
        "image/x-emf": "emf",  # Define a MIME mapping.
        "image/x-wmf": "wmf",  # Define a MIME mapping.
    }
    return mapping.get((content_type or "").lower(), "png")  # Return the computed value.

def _normalise_image_blob(blob: bytes, extension: str, filename: str):  # Define the image normaliser.
    """Normalise an image blob."""
    extension = (extension or "png").lower()  # Store the detected extension.

    if extension in {"emf", "wmf"}:  # Check the current condition.
        return blob, extension, filename  # Return the computed value.

    try:  # Start guarded parsing.
        with Image.open(io.BytesIO(blob)) as img:  # Open the resource safely.
            output = io.BytesIO()  # Buffer the converted image.

            save_format = "PNG" if img.mode in {"RGBA", "LA", "P"} else "JPEG"  # Choose the image save format.
            converted = img.convert("RGBA" if save_format == "PNG" else "RGB")  # Convert the image mode.
            converted.save(output, format=save_format)  # Write the converted image.

            new_extension = "png" if save_format == "PNG" else "jpg"  # Store the new extension.
            root, _ = os.path.splitext(filename)  # Split the filename root.
            return output.getvalue(), new_extension, f"{root}.{new_extension}"  # Return the computed value.

    except (UnidentifiedImageError, OSError):  # Handle parsing failures.
        return blob, extension, filename  # Return the computed value.
